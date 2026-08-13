"""MechVault services: auth, AI, extraction, watermark."""
import os, io, httpx, fitz, threading
from itsdangerous import URLSafeTimedSerializer
from argon2 import PasswordHasher
from . import db_models

SECRET = os.environ.get("MECHVAULT_SECRET", "dev-secret-change-me")
SIG = URLSafeTimedSerializer(SECRET)
PH = PasswordHasher()

# ---------- auth ----------
def hash_password(p): return PH.hash(p)
def verify_password(p, h):
    try: return PH.verify(h, p)
    except Exception: return False

def make_session(email): return SIG.dumps(email)
def read_session(token, max_age=60 * 60 * 24 * 30):
    try: return SIG.loads(token, max_age=max_age)
    except Exception: return None

def get_user_by_email(email):
    db = db_models.SessionLocal()
    try: return db.query(db_models.User).filter_by(email=email).first()
    finally: db.close()


# ---------- AI (opencode-zen, free tier) ----------
OPENCODE_ZEN_KEY = os.environ.get("OPENCODE_ZEN_API_KEY")
OPENCODE_ZEN_URL = "https://opencode.ai/zen/v1/chat/completions"
AI_MODEL = os.environ.get("MECHVAULT_AI_MODEL", "nemotron-3-ultra-free")


def ai_available():
    return bool(OPENCODE_ZEN_KEY)


def ai_generate(system, prompt, model=None, max_tokens=4000):
    if not OPENCODE_ZEN_KEY:
        raise RuntimeError("OPENCODE_ZEN_API_KEY not set — AI generation disabled.")
    model = model or AI_MODEL
    resp = httpx.post(OPENCODE_ZEN_URL,
                      headers={"Authorization": f"Bearer {OPENCODE_ZEN_KEY}",
                               "Content-Type": "application/json"},
                      json={"model": model,
                            "messages": [{"role": "system", "content": system},
                                         {"role": "user", "content": prompt}],
                            "max_tokens": max_tokens, "temperature": 0.3},
                      timeout=180)
    resp.raise_for_status()
    return resp.json()["choices"][0]["message"]["content"]


def build_note_prompt(topic, module_name, subject_name, sources):
    sys = ("You are a top Mechanical Engineering tutor. Produce rigorous, exam-focused "
           "study notes in GitHub-flavored Markdown. Use sections: "
           "## Concept, ## Key Formulas, ## Derivation/Method, ## Worked Example, "
           "## Common Pitfalls, ## Exam Tips. Be precise with equations.")
    body = f"Subject: {subject_name}\nModule: {module_name}\nTopic: {topic}\n\n"
    body += "SOURCE MATERIAL (extract the important parts):\n"
    for label, text in sources:
        if text and text.strip():
            body += f"\n### {label}\n{text[:6000]}\n"
    return sys, body


# ---------- extraction ----------
def extract_pdf(path):
    doc = fitz.open(path)
    try:
        return "\n".join(page.get_text() for page in doc)
    finally:
        doc.close()


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return "\n".join(out)


def add_watermark(src_path, email):
    """Return watermarked PDF bytes (faint email on every page)."""
    doc = fitz.open(src_path)
    for page in doc:
        r = page.rect
        page.insert_text((r.width / 2 - 150, r.height / 2), f"{email}  ·  MechVault",
                         fontsize=22, color=(0.82, 0.82, 0.82), rotate=0, overlay=True)
        page.insert_text((16, r.height - 16), f"Confidential — {email}",
                         fontsize=8, color=(0.6, 0.6, 0.6), rotate=0, overlay=True)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


# ---------- background generation ----------
def run_generation(job_id, topic_id):
    db = db_models.SessionLocal()
    try:
        job = db.query(db_models.Job).get(job_id)
        job.status = "running"; db.commit()
        topic = db.query(db_models.Topic).get(topic_id)
        module = db.query(db_models.Module).get(topic.module_id)
        subject = db.query(db_models.Subject).get(module.subject_id)
        resources = db.query(db_models.Resource).filter_by(topic_id=topic_id).all()
        papers = db.query(db_models.Paper).filter_by(topic_id=topic_id).all()
        sources = []
        for r in resources:
            sources.append((f"{r.kind}: {r.filename}", r.extracted_text))
        for p in papers:
            sources.append((f"{p.kind} paper: {p.title}", p.extracted_text))
        sys_p, prompt = build_note_prompt(topic.name, module.name, subject.name, sources)
        markdown = ai_generate(sys_p, prompt)
        note = db.query(db_models.Note).filter_by(topic_id=topic_id).first()
        if not note:
            note = db_models.Note(topic_id=topic_id)
            db.add(note)
        note.content_markdown = markdown
        note.status = "draft"
        note.model = AI_MODEL
        db.commit()
        job.status = "done"; job.message = "Notes generated (draft)."
        db.commit()
    except Exception as e:
        db.rollback()
        job = db.query(db_models.Job).get(job_id)
        job.status = "error"; job.message = str(e)[:500]; db.commit()
    finally:
        db.close()


def start_generation(topic_id):
    db = db_models.SessionLocal()
    job = db_models.Job(kind="generate_note", ref_id=topic_id, status="pending")
    db.add(job); db.commit(); jid = job.id; db.close()
    threading.Thread(target=run_generation, args=(jid, topic_id), daemon=True).start()
    return jid
